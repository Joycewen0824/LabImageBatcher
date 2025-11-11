#!/usr/bin/env python3
"""
Lab Image Batcher — 实验室显微/蛋白/细胞图片批处理与排版小应用（Streamlit）

功能：
1) 批量统一尺寸（按最长边/短边/指定宽高，支持等比、禁止放大、插值算法选择）。
2) 批量补边（pad）或裁剪（crop），生成统一画布大小，背景可选（白/黑/自定义）。
3) 批量重命名与导出（ZIP）。
4) 生成排版拼图（contact sheet/蒙太奇）：按列数或行列指定，间距、边距、画布尺寸（像素或A4/Letter @DPI），文件名字幕可选。
5) 元数据表（CSV）：原始尺寸、目标尺寸、缩放比例、输出文件名。
6) 可选：将每张图导出到 PPTX（一页一图，标题=文件名，可选居中与标注尺寸）。

运行：
    pip install streamlit pillow numpy python-pptx
    streamlit run app.py

文件名：建议保存为 app.py

注意：
- 显微图的尺度条：本工具只做像素层面的缩放。如果需要保持物理尺度，请使用“禁止放大”与“等比缩放”，并在拼图时关闭再次缩放。
- TIFF 支持：Pillow 对部分多通道/多页 tiff 支持有限，此处取第一页；如需拓展可自行修改。
"""

from __future__ import annotations
import io
import math
import zipfile
from dataclasses import dataclass
from typing import List, Tuple, Optional

import numpy as np
from PIL import Image, ImageOps, ImageDraw, ImageFont
import streamlit as st

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# ---------------------------- 工具函数 ----------------------------
INTERP_MAP = {
    "最近邻 (Nearest)": Image.NEAREST,
    "双线性 (Bilinear)": Image.BILINEAR,
    "双三次 (Bicubic)": Image.BICUBIC,
    "Lanczos (高质量)": Image.LANCZOS,
}

DEFAULT_FONT = None  # 让 Pillow 使用默认字体；如需中文字体，可改为本地 .ttf 路径


def load_image(file) -> Image.Image:
    img = Image.open(file)
    if getattr(img, "n_frames", 1) > 1:
        img.seek(0)
    return img.convert("RGB")


def parse_size(text: str) -> Tuple[int, int]:
    """解析类似 "1024x768" 的输入为 (w,h)。"""
    parts = text.lower().replace("×", "x").split("x")
    if len(parts) != 2:
        raise ValueError("输入格式应为 宽x高，例如 1024x768")
    w, h = int(parts[0].strip()), int(parts[1].strip())
    if w <= 0 or h <= 0:
        raise ValueError("宽高需要是正整数")
    return w, h


def pad_to_size(img: Image.Image, target_size: Tuple[int, int], color=(255, 255, 255)) -> Image.Image:
    tw, th = target_size
    return ImageOps.pad(img, target_size, color=color, centering=(0.5, 0.5)) if False else ImageOps.expand(
        ImageOps.fit(img, (min(img.width, tw), min(img.height, th)), method=Image.NEAREST, centering=(0.5, 0.5)),
        border=(max(0, (tw - min(img.width, tw)) // 2), max(0, (th - min(img.height, th)) // 2)), fill=color
    )


def resize_by_long_or_short(img: Image.Image, target: int, mode: str, keep_ratio: bool, no_upscale: bool, interp) -> Image.Image:
    w, h = img.size
    if keep_ratio:
        if mode == "最长边":
            scale = target / max(w, h)
        else:
            scale = target / min(w, h)
        if no_upscale and scale > 1.0:
            return img
        nw, nh = max(1, int(round(w * scale))), max(1, int(round(h * scale)))
    else:
        if mode == "最长边":
            if w >= h:
                nw, nh = target, h
            else:
                nw, nh = w, target
        else:  # 短边
            if w <= h:
                nw, nh = target, h
            else:
                nw, nh = w, target
        if no_upscale:
            nw, nh = min(nw, w), min(nh, h)
    return img.resize((nw, nh), interp)


def resize_to_box(img: Image.Image, tw: int, th: int, fit_mode: str, color, interp, no_upscale: bool) -> Image.Image:
    w, h = img.size
    if fit_mode == "等比缩放，补边":
        scale = min(tw / w, th / h)
        if no_upscale:
            scale = min(1.0, scale)
        nw, nh = max(1, int(round(w * scale))), max(1, int(round(h * scale)))
        img2 = img.resize((nw, nh), interp)
        canvas = Image.new("RGB", (tw, th), color)
        canvas.paste(img2, ((tw - nw) // 2, (th - nh) // 2))
        return canvas
    elif fit_mode == "等比填满，居中裁剪":
        scale = max(tw / w, th / h)
        if no_upscale:
            scale = min(1.0, scale)
        nw, nh = max(1, int(round(w * scale))), max(1, int(round(h * scale)))
        img2 = img.resize((nw, nh), interp)
        left = max(0, (nw - tw) // 2)
        top = max(0, (nh - th) // 2)
        return img2.crop((left, top, left + tw, top + th))
    else:  # 直接拉伸
        if no_upscale:
            tw, th = min(tw, w), min(th, h)
        return img.resize((tw, th), interp)


def draw_caption(img: Image.Image, text: str, font_size: int, color=(0, 0, 0), bg: Optional[Tuple[int,int,int]] = None, pad: int = 6) -> Image.Image:
    if not text:
        return img
    try:
        font = ImageFont.truetype(DEFAULT_FONT, font_size) if DEFAULT_FONT else ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()
    draw = ImageDraw.Draw(img)
    tw, th = draw.textbbox((0, 0), text, font=font)[2:]
    W, H = img.size
    cap_h = th + pad * 2
    canvas = Image.new("RGB", (W, H + cap_h), (255, 255, 255) if bg is None else bg)
    canvas.paste(img, (0, 0))
    draw = ImageDraw.Draw(canvas)
    tx = (W - tw) // 2
    ty = H + pad
    draw.text((tx, ty), text, fill=color, font=font)
    return canvas


@dataclass
class ProcessedImage:
    name: str
    orig_size: Tuple[int, int]
    out_img: Image.Image
    out_size: Tuple[int, int]
    scale: float


# ---------------------------- Streamlit UI ----------------------------
st.set_page_config(page_title="Lab Image Batcher", layout="wide")
st.title("🧪 Lab Image Batcher｜显微/蛋白/细胞图 批量处理与排版")

with st.sidebar:
    st.header("① 导入图片")
    files = st.file_uploader("选择图片（支持多选：JPG/PNG/TIF/TIFF/BMP）", type=["jpg","jpeg","png","tif","tiff","bmp"], accept_multiple_files=True)
    st.caption("注：TIFF 仅取第一页；所有图片将转为 RGB。")

    st.header("② 尺寸与缩放")
    mode = st.selectbox("方式", ["按最长/短边", "指定画布 (宽x高)"])
    keep_ratio = True
    no_upscale = st.checkbox("禁止放大", value=True)
    interp_name = st.selectbox("插值算法", list(INTERP_MAP.keys()), index=3)
    interp = INTERP_MAP[interp_name]

    pad_bg = st.color_picker("补边/画布背景色", value="#FFFFFF")
    bg_rgb = tuple(int(pad_bg.lstrip('#')[i:i+2], 16) for i in (0,2,4))

    if mode == "按最长/短边":
        ls_mode = st.radio("基准边", ["最长边", "短边"], horizontal=True)
        target_px = st.number_input("目标像素", min_value=1, value=1024)
    else:
        box_str = st.text_input("目标画布，例如 1024x768", value="1024x768")
        fit_mode = st.selectbox("适配方式", ["等比缩放，补边", "等比填满，居中裁剪", "不保持比例，直接拉伸"])
        try:
            tw, th = parse_size(box_str)
        except Exception as e:
            st.error(str(e))
            tw, th = 1024, 768

    st.header("③ 拼图/排版")
    make_contact = st.checkbox("生成拼图 (contact sheet)", value=True)
    cols = st.number_input("每行列数", min_value=1, value=4)
    grid_gap = st.number_input("格子间距 (px)", min_value=0, value=12)
    margin = st.number_input("画布边距 (px)", min_value=0, value=24)
    show_caption = st.checkbox("添加文件名字幕", value=False)
    caption_font = st.number_input("字幕字号", min_value=6, value=14)

    sheet_size_mode = st.radio("画布尺寸", ["自动按网格", "自定义像素", "A4/Letter+DPI"], index=0)
    if sheet_size_mode == "自定义像素":
        sheet_wh = st.text_input("画布宽x高 (px)", value="2480x3508")
        try:
            sheet_w, sheet_h = parse_size(sheet_wh)
        except Exception:
            sheet_w, sheet_h = 2480, 3508
    elif sheet_size_mode == "A4/Letter+DPI":
        paper = st.selectbox("纸张", ["A4", "Letter"])
        dpi = st.number_input("DPI", min_value=72, value=300)
        if paper == "A4":
            sheet_w, sheet_h = int(8.27 * dpi), int(11.69 * dpi)  # 竖向
        else:
            sheet_w, sheet_h = int(8.5 * dpi), int(11 * dpi)
    else:
        sheet_w, sheet_h = None, None

    st.header("④ 导出")
    do_zip = st.checkbox("导出 ZIP (单张处理后)", value=True)
    do_csv = st.checkbox("导出 CSV 元数据", value=True)
    do_pptx = st.checkbox("导出 PPTX (一页一图)", value=False and PPTX_AVAILABLE)
    if do_pptx and not PPTX_AVAILABLE:
        st.warning("未安装 python-pptx，无法导出 PPTX。请先 pip install python-pptx")

# ---------------------------- 处理逻辑 ----------------------------
processed: List[ProcessedImage] = []

if files:
    for f in files:
        try:
            img = load_image(f)
        except Exception as e:
            st.error(f"读取失败：{f.name} — {e}")
            continue
        orig = img.size
        if mode == "按最长/短边":
            out = resize_by_long_or_short(img, target_px, ls_mode, keep_ratio=True, no_upscale=no_upscale, interp=interp)
        else:
            out = resize_to_box(img, tw, th, fit_mode, bg_rgb, interp, no_upscale=no_upscale)
        processed.append(ProcessedImage(name=f.name, orig_size=orig, out_img=out, out_size=out.size, scale=min(out.size[0]/orig[0], out.size[1]/orig[1])))

    st.success(f"已处理 {len(processed)} 张图像")

    # 预览
    st.subheader("预览")
    preview_cols = st.columns(4)
    for i, item in enumerate(processed[:8]):
        with preview_cols[i % 4]:
            st.image(item.out_img, caption=f"{item.name} → {item.out_size[0]}x{item.out_size[1]}", use_column_width=True)

    # ---------------- ZIP 导出 ----------------
    if do_zip:
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
            for idx, item in enumerate(processed, 1):
                base = item.name.rsplit('.', 1)[0]
                out_name = f"{idx:03d}_{base}.jpg"
                img_bytes = io.BytesIO()
                item.out_img.save(img_bytes, format="JPEG", quality=95)
                zf.writestr(out_name, img_bytes.getvalue())
        zip_buf.seek(0)
        st.download_button("⬇️ 下载处理后的图片 (ZIP)", data=zip_buf, file_name="processed_images.zip", mime="application/zip")

    # ---------------- CSV 元数据 ----------------
    if do_csv:
        import csv
        csv_buf = io.StringIO()
        writer = csv.writer(csv_buf)
        writer.writerow(["filename", "orig_w", "orig_h", "out_w", "out_h", "scale"])
        for item in processed:
            writer.writerow([item.name, item.orig_size[0], item.orig_size[1], item.out_size[0], item.out_size[1], f"{item.scale:.4f}"])
        st.download_button("⬇️ 下载元数据 (CSV)", data=csv_buf.getvalue().encode("utf-8-sig"), file_name="image_metadata.csv", mime="text/csv")

    # ---------------- PPTX 导出 ----------------
    if do_pptx and PPTX_AVAILABLE:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for item in processed:
            slide = prs.slides.add_slide(blank)
            # 以 10x7.5 英寸内容区估算（默认宽 13.333" 高 7.5"，此处简单铺满高）
            pic_stream = io.BytesIO()
            item.out_img.save(pic_stream, format="PNG")
            pic_stream.seek(0)
            slide.shapes.add_picture(pic_stream, Inches(1), Inches(1), width=Inches(8))
        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
        st.download_button("⬇️ 下载 PPTX", data=pptx_buf, file_name="images.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # ---------------- 拼图/排版 ----------------
    if make_contact:
        # 先可选地给每张图加字幕
        imgs_for_grid = []
        for item in processed:
            im = item.out_img
            if show_caption:
                im = draw_caption(im, item.name, caption_font, color=(0,0,0), bg=(255,255,255))
            imgs_for_grid.append(im)

        # 计算格子尺寸（使用当前处理后图片的最大宽/高）
        cell_w = max(im.width for im in imgs_for_grid)
        cell_h = max(im.height for im in imgs_for_grid)

        if sheet_w is None or sheet_h is None:  # 自动根据列数排版
            rows = math.ceil(len(imgs_for_grid) / cols)
            W = margin*2 + cols*cell_w + (cols-1)*grid_gap
            H = margin*2 + rows*cell_h + (rows-1)*grid_gap
        else:
            W, H = sheet_w, sheet_h
            # 如果自定义画布过小，提示
            minW = margin*2 + cols*cell_w + (cols-1)*grid_gap
            rows = max(1, math.ceil((len(imgs_for_grid))/cols))
            minH = margin*2 + rows*cell_h + (rows-1)*grid_gap
            if W < minW or H < minH:
                st.warning("自定义画布可能过小，部分图片会被截断或超出画布。")

        sheet = Image.new("RGB", (W, H), bg_rgb)

        # 布局粘贴
        for i, im in enumerate(imgs_for_grid):
            r = i // cols
            c = i % cols
            x = margin + c*(cell_w + grid_gap)
            y = margin + r*(cell_h + grid_gap)
            # 如果图片比 cell 小，则居中；否则左上对齐粘贴（已统一尺寸的通常不会溢出）
            ox = x + (cell_w - im.width)//2
            oy = y + (cell_h - im.height)//2
            if ox < 0 or oy < 0:
                ox, oy = x, y
            sheet.paste(im, (ox, oy))

        # 预览与导出
        st.image(sheet, caption=f"拼图预览：{W}x{H}", use_column_width=True)
        out_png = io.BytesIO()
        sheet.save(out_png, format="PNG")
        out_png.seek(0)
        st.download_button("⬇️ 下载拼图 (PNG)", data=out_png, file_name="contact_sheet.png", mime="image/png")

else:
    st.info("请在左侧选择要处理的图片。")
